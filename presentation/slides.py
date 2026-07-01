#!/usr/bin/env python3
"""Slide content for the RIOS product-walkthrough deck.

Grounded in what the application actually does. RIOS today is a FOUNDATION /
vertical slice — the reinsurance core (place → bind → account → reconcile → claims)
is proven end to end; enterprise breadth is designed-for and named honestly, per
docs/phases.md and docs/open-questions.md. No invented metrics; no black-box AI
claims (the assistant is a deterministic, grounded, confirmation-gated engine)."""
from build_deck import *  # noqa
from build_deck import _noshadow, _mix
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

N = [0]
def pn():
    N[0] += 1
    return N[0]


def chip_flow(s, x, y, maxw, items, color, size=9.5, ch=0.32, gap=0.14):
    """Flow pill chips, wrapping to new rows. Each pill is sized to fit its label
    (measured width + padding) so text never wraps inside the pill."""
    pad = 0.32
    cx, cy = x, y
    for it in items:
        w = text_w(it, size, True) + pad
        if cx > x and cx + w > x + maxw + 0.01:
            cx = x; cy += ch + gap
        rect(s, cx, cy, w, ch, fill=_mix(color, WHITE, 0.12), radius=0.5)
        tb = text(s, cx, cy - 0.02, w, ch, [[(it, size, _mix(color, INK, 0.62), True)]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tb.text_frame.word_wrap = False  # never wrap inside a pill
        cx += w + gap
    return cy + ch


def label(s, x, y, txt, color):
    text(s, x, y, 5.6, 0.3, [[(txt, 10.5, color, True)]])


# ============================================================== 1. COVER
def cover():
    s = slide(WHITE)
    rect(s, 0, 0, 5.6, 7.5, fill=INK, rounded=False)
    accent_bar(s, 0, 0, 0.16, 7.5, BLUE)
    m = s.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(0.9), Inches(0.9), Inches(0.6), Inches(0.6))
    _noshadow(m); m.fill.solid(); m.fill.fore_color.rgb = BLUE; m.line.fill.background()
    text(s, 1.62, 0.92, 3.5, 0.6, [[('RIOS', 24, WHITE, True)]])
    text(s, 1.62, 1.35, 4, 0.3, [[('Reinsurance Intelligent OS', 10, MUTE, False)]])
    text(s, 0.9, 2.55, 4.5, 2.2, [
        [('Reinsurance', 33, WHITE, True)],
        [('Intelligence &', 33, WHITE, True)],
        [('Operations Suite', 33, RGBColor(0x93,0xC5,0xFD), True)],
    ], line_spacing=1.03)
    text(s, 0.9, 4.85, 4.4, 1.1, [[('A unified operating system for reinsurance — placement, underwriting, accounting, claims, analytics and the back office on one governed platform.', 12.5, RGBColor(0xC7,0xD2,0xFE), False)]], line_spacing=1.24)
    text(s, 0.9, 6.6, 4.4, 0.4, [[('Product Walkthrough', 12, WHITE, True)]])
    text(s, 6.1, 0.95, 6.6, 0.4, [[('ENTERPRISE PLATFORM FOR MODERN REINSURANCE', 11, BLUE, True)]])
    browser_frame(s, 'executive.png', 6.1, 1.5, 6.5, addr='app.rios.cloud/executive')
    text(s, 6.1, 6.68, 6.6, 0.4, [[('Place   ·   Bind   ·   Account   ·   Reconcile   ·   Claims  —  on one platform.', 12, SLATE, True)]])

# ============================================================== 2. ABOUT
def about():
    s = slide(BG); accent_bar(s); kicker(s, 'About RIOS'); title(s, 'One platform. One data model. One source of truth.')
    text(s, 0.9, 1.72, 11.4, 1.3, [[('RIOS unifies the reinsurance value chain — placement, underwriting, treaty & facultative administration, technical and general accounting, claims, analytics and the back office — into a single, metadata-driven, multi-tenant system, replacing the disconnected legacy tools and spreadsheets most reinsurers run today.', 13.5, SLATE, False)]], line_spacing=1.25)
    cards = [('What it is', 'A foundation that proves the reinsurance core end to end — place → bind → account → reconcile → claims — correct, secure and audited.', BLUE),
             ('How it is built', 'Metadata-driven configuration, integer-accurate money, hash-chained audit and role-based access — foundations, not add-ons.', INDIGO),
             ('Why it matters', 'The technical → financial chain reconciles to zero and every material change is audited — correctness you can inspect.', GREEN)]
    for i, (h, b, c) in enumerate(cards):
        x = 0.9 + i*4.05
        rect(s, x, 3.35, 3.8, 3.15, fill=WHITE, line=LINE, radius=0.06)
        icon_tile(s, x+0.35, 3.72, color=c, size=0.6)
        text(s, x+0.35, 4.5, 3.1, 0.4, [[(h, 16, INK, True)]])
        text(s, x+0.35, 4.95, 3.15, 1.5, [[(b, 11.5, SLATE, False)]], line_spacing=1.2)
    page_no(s, pn())

# ============================================================== 3. PROBLEMS
def problems():
    s = slide(WHITE); accent_bar(s, color=AMBER); kicker(s, 'The Problem', color=AMBER)
    title(s, 'Reinsurance still runs on disconnected systems')
    items = [('Manual underwriting', 'Excel and email drive placement and pricing.'),
             ('Siloed legacy suites', 'Policy, claims and finance never share one truth.'),
             ('Duplicate data entry', 'The same risk rekeyed across many systems.'),
             ('Slow placement & pricing', 'Time lost to spreadsheets and version chaos.'),
             ('No real-time visibility', 'Exposure and capacity are always yesterday’s.'),
             ('Fragmented reporting', 'Every return stitched together by hand.'),
             ('Regulatory complexity', 'Solvency, IFRS and audit run offline.'),
             ('No embedded assistance', 'No decision support where underwriters work.')]
    for i, (h, b) in enumerate(items):
        col = i % 4; row = i // 4
        x = 0.9 + col*3.05; y = 1.95 + row*2.3
        rect(s, x, y, 2.85, 2.0, fill=BG, line=LINE, radius=0.07)
        icon_tile(s, x+0.28, y+0.28, color=AMBER, size=0.48)
        text(s, x+0.28, y+0.92, 2.4, 0.4, [[(h, 13, INK, True)]])
        text(s, x+0.28, y+1.3, 2.45, 0.7, [[(b, 10.5, SLATE, False)]], line_spacing=1.12)
    page_no(s, pn())

# ============================================================== 4. PRODUCT MAP
def product_map():
    s = slide(BG); accent_bar(s); kicker(s, 'Product Map'); title(s, 'Every module, one connected suite', size=27)
    groups = [
        ('Overview', ['Dashboard','Executive','Intelligence','AI Insights','Search','Mobile'], BLUE),
        ('Underwriting', ['Underwriting Workspace','Treaty','Facultative','Placement','Pricing','Capacity & Exposure','Territory','Retrocession','Adjustments'], INDIGO),
        ('Distribution', ['Parties','Clients','Brokers','Cedents','CRM'], GREEN),
        ('Operations', ['Claims','Bordereaux','Recoveries','Operations Center','Workflow Center','Audit Log'], AMBER),
        ('Finance', ['Accounting','Statements','Finance','Treasury','Period Close','Procurement'], BLUE),
        ('Analytics & Compliance', ['Reports','Scheduled Reports','Analytics','Risk & Capital','Regulatory','Compliance','Returns'], INDIGO),
        ('HRMS', ['Attendance','People','Payroll','Performance','Assets','Org Structure'], GREEN),
        ('Master Data', ['Products','Reference Data','Code Lists'], AMBER),
        ('Documents & Knowledge', ['Documents','Templates','Knowledge Base'], BLUE),
        ('Integration & Automation', ['Integration Hub','Messaging','Automation Studio','Portal'], INDIGO),
        ('Administration', ['Admin','Legal Entities','Ops Console','Delegation','Security','Security Ops','Field Security','Retention','Cost Mgmt','Features'], GREEN),
        ('Ask RIOS Assistant', ['Copilot','Grounded Search','Recommendations'], AMBER)]
    for i, (h, items, c) in enumerate(groups):
        col = i % 4; row = i // 4
        x = 0.9 + col*3.05; y = 1.75 + row*1.75
        rect(s, x, y, 2.85, 1.58, fill=WHITE, line=LINE, radius=0.07)
        rect(s, x, y, 2.85, 0.1, fill=c, rounded=False)
        text(s, x+0.22, y+0.17, 2.5, 0.35, [[(h, 12, INK, True)]])
        text(s, x+0.22, y+0.56, 2.5, 1.0, [[(' · '.join(items), 8.3, SLATE, False)]], line_spacing=1.12)
    page_no(s, pn())

# ============================================================== 5. ARCHITECTURE
def architecture():
    s = slide(WHITE); accent_bar(s, color=INDIGO); kicker(s, 'System Architecture', color=INDIGO)
    title(s, 'A modern, layered enterprise architecture')
    layers = [('Experience','Web · Mobile · Portals · RIOS Assistant', BLUE),
              ('Business Modules','Underwriting · Distribution · Operations · Finance · HRMS · Analytics', INDIGO),
              ('Platform Services','Workflow Engine · Rules · Notifications · Audit · Search · AI', GREEN),
              ('Data & Integration','PostgreSQL (Row-Level Security) · Outbox + Relay · APIs · ERP / Email / DMS', AMBER)]
    y = 2.0
    for i, (h, b, c) in enumerate(layers):
        rect(s, 1.4, y, 8.4, 1.05, fill=_mix(c, WHITE, 0.10), line=None, radius=0.05)
        rect(s, 1.4, y, 0.14, 1.05, fill=c, rounded=False)
        text(s, 1.75, y+0.14, 3, 0.4, [[(h, 15, INK, True)]])
        text(s, 1.75, y+0.56, 7.9, 0.4, [[(b, 11.5, SLATE, False)]])
        if i < 3:
            a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(5.4), Inches(y+1.02), Inches(0.4), Inches(0.2))
            _noshadow(a); a.fill.solid(); a.fill.fore_color.rgb = MUTE; a.line.fill.background()
        y += 1.28
    rect(s, 10.2, 2.0, 2.5, 4.35, fill=INK, radius=0.06)
    text(s, 10.45, 2.25, 2.1, 0.4, [[('PLATFORM', 11, RGBColor(0x93,0xC5,0xFD), True)]])
    for i, t in enumerate(['Multi-tenant + RLS','Zero-trust access','Hash-chained audit','Integer-money ledger','Modular monolith','Container-based deploy']):
        text(s, 10.45, 2.75+i*0.58, 2.1, 0.4, [[('•  '+t, 11, WHITE, False)]])
    page_no(s, pn())

# ============================================================== 6. SECURITY
def security():
    s = slide(BG); accent_bar(s); kicker(s, 'Security & Trust'); title(s, 'Secure and auditable by design')
    items = [('Row-Level Security','Postgres RLS isolates every tenant — fail-closed.'),
             ('RBAC + permissions','Role-based, permission-bound; admin override is explicit.'),
             ('OIDC SSO + MFA','OIDC wired; TOTP / passkey ceremonies built; SAML at deployment.'),
             ('KMS encryption','Envelope encryption (AES-256-GCM); managed HSM/KMS in prod.'),
             ('Tamper-evident audit','Hash-chained, append-only audit trail on mutations.'),
             ('Least privilege','Separate owner vs application DB roles; no back doors.'),
             ('Guardrailed AI','The assistant re-checks permissions and confirms before acting.'),
             ('Control framework','Designed around GDPR / ISO 27001 / SOC 2 principles.')]
    for i, (h, b) in enumerate(items):
        col = i % 4; row = i // 4
        x = 0.9 + col*3.05; y = 1.95 + row*2.3
        rect(s, x, y, 2.85, 2.0, fill=WHITE, line=LINE, radius=0.07)
        icon_tile(s, x+0.28, y+0.26, color=GREEN, size=0.48)
        text(s, x+0.28, y+0.9, 2.4, 0.4, [[(h, 12.5, INK, True)]])
        text(s, x+0.28, y+1.28, 2.5, 0.7, [[(b, 10.0, SLATE, False)]], line_spacing=1.12)
    page_no(s, pn())

# ============================================================== SECTION DIVIDER
def section(num, ttl, sub, color=BLUE):
    s = slide(INK)
    accent_bar(s, 0, 0, 0.16, 7.5, color)
    text(s, 0.9, 2.5, 3, 1.2, [[(num, 60, _mix(color, INK, 0.9), True)]])
    text(s, 0.9, 3.55, 11.5, 1.0, [[(ttl, 33, WHITE, True)]])
    rect(s, 0.95, 4.42, 0.8, 0.05, fill=color, rounded=False)
    text(s, 0.95, 4.58, 10.5, 0.8, [[(sub, 14, RGBColor(0xC7,0xD2,0xFE), False)]], line_spacing=1.2)
    return s

# ============================================================== MODULE SLIDE
def module_slide(kick, ttl, shot, purpose, features, inside, users, ai, addr, color=BLUE):
    s = slide(WHITE); accent_bar(s, color=color); kicker(s, kick, y=0.82, color=color)
    text(s, 0.9, 1.1, 5.6, 0.7, [[(ttl, 22, INK, True)]], line_spacing=0.98)
    rect(s, 0.92, 1.78, 0.62, 0.05, fill=color, rounded=False)
    text(s, 0.9, 1.96, 5.55, 1.15, [[(purpose, 11.5, SLATE, False)]], line_spacing=1.2)
    # screenshot
    browser_frame(s, shot, 6.6, 1.1, 6.0, addr=addr)
    # capabilities
    label(s, 0.9, 3.28, 'KEY CAPABILITIES', color)
    bullets(s, 0.9, 3.62, 5.55, 1.3, features, size=11, marker_color=color, gap=4)
    # what's inside
    label(s, 0.9, 4.95, 'WHAT’S INSIDE', INK)
    ib = chip_flow(s, 0.9, 5.28, 5.55, inside, color)
    # primary users — placed below wherever the inside chips actually end
    uy = max(6.10, ib + 0.18)
    label(s, 0.9, uy, 'PRIMARY USERS', MUTE)
    chip_flow(s, 0.9, uy + 0.30, 5.55, users, SLATE, size=9)
    # AI strip
    rect(s, 6.6, 5.25, 6.0, 1.35, fill=_mix(INDIGO, WHITE, 0.07), radius=0.07)
    ai_badge(s, 6.82, 5.45)
    text(s, 7.44, 5.5, 4.9, 0.3, [[('INTELLIGENCE IN THIS MODULE', 9.5, INDIGO, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.85, 5.9, 5.55, 0.65, [[(ai, 10.3, SLATE, False)]], line_spacing=1.1)
    page_no(s, pn())

# ============================================================== ASSISTANT
def assistant():
    s = slide(WHITE); accent_bar(s, color=INDIGO); kicker(s, 'RIOS Assistant', y=0.82, color=INDIGO)
    text(s, 0.9, 1.1, 5.6, 0.7, [[('Ask RIOS — a grounded, guardrailed assistant', 21, INK, True)]], line_spacing=1.0)
    rect(s, 0.92, 1.86, 0.62, 0.05, fill=INDIGO, rounded=False)
    text(s, 0.9, 2.04, 5.55, 1.5, [[('A deterministic intent engine — not a black-box LLM. It answers from your own tenant data, can prepare an action, and always re-checks permissions and asks for explicit confirmation before it changes anything. The whole platform works fully with AI switched off.', 11.5, SLATE, False)]], line_spacing=1.22)
    browser_frame(s, 'ai-insights.png', 6.6, 1.1, 6.0, addr='app.rios.cloud/ai-insights')
    label(s, 0.9, 3.75, 'HOW IT WORKS', INDIGO)
    bullets(s, 0.9, 4.1, 5.55, 2.2, [
        'Grounded queries across every module',
        'Answers explained from your data — no black box',
        'Prepares actions, then confirms before committing',
        'Re-checks permissions server-side — no backdoor',
        'Optional LLM only narrates — it never invents numbers',
    ], size=11.5, marker_color=INDIGO, gap=6)
    rect(s, 6.6, 5.25, 6.0, 1.35, fill=_mix(GREEN, WHITE, 0.08), radius=0.07)
    text(s, 6.85, 5.4, 5.5, 1.1, [[('Guardrail: every mutating action is confirmation-gated and permission-checked. The assistant is a copilot, never an autopilot — and it never invents your numbers.', 11, SLATE, False)]], line_spacing=1.15)
    page_no(s, pn())

# ============================================================== PRODUCT LIFECYCLE
def lifecycle():
    s = slide(BG); accent_bar(s, color=GREEN); kicker(s, 'Product Lifecycle', color=GREEN)
    title(s, 'The reinsurance lifecycle, on one platform')
    steps = [('Place','Submission, quotes, placement tower', BLUE),
             ('Bind','Terms agreed; deposit premium booked', INDIGO),
             ('Account','Statements & balanced GL postings', AMBER),
             ('Reconcile','Technical → financial chain to zero', GREEN),
             ('Claims','FNOL, reserves, recoveries, cash calls', BLUE)]
    y = 2.6
    for i, (t, b, c) in enumerate(steps):
        x = 0.9 + i*2.42
        rect(s, x, y, 2.15, 2.0, fill=WHITE, line=LINE, radius=0.09)
        rect(s, x, y, 2.15, 0.12, fill=c, rounded=False)
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.24), Inches(y+0.34), Inches(0.5), Inches(0.5))
        _noshadow(d); d.fill.solid(); d.fill.fore_color.rgb = c; d.line.fill.background()
        text(s, x+0.24, y+0.34, 0.5, 0.5, [[(str(i+1), 16, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x+0.2, y+1.02, 1.85, 0.35, [[(t, 14, INK, True)]])
        text(s, x+0.2, y+1.38, 1.85, 0.6, [[(b, 9.5, SLATE, False)]], line_spacing=1.1)
        if i < len(steps)-1:
            arrow(s, x+2.16, y+0.85, 0.26, 0.28, MUTE)
    text(s, 0.9, 5.2, 11.5, 0.9, [[('This is the delivered vertical slice — integration-tested: a bound catastrophe XL treaty flows from placement through binding, statement, balanced general-ledger postings and reconciliation to zero, into claims. Every step is governed, audited and reconcilable.', 12.5, SLATE, False)]], line_spacing=1.25)
    page_no(s, pn())

# ============================================================== TECHNOLOGY
def technology():
    s = slide(WHITE); accent_bar(s, color=INDIGO); kicker(s, 'Technology', color=INDIGO)
    title(s, 'Built on a modern, proven stack')
    cols = [('Experience', ['React + Vite','Design-token system','Responsive & mobile','Accessible UI']),
            ('Services', ['Fastify (Node / TS)','Pure domain core','Workflow & rules engine','Transactional outbox']),
            ('Data', ['PostgreSQL 16','Row-Level Security','Integer-money ledger','Hash-chained audit']),
            ('Platform', ['Docker Compose','CI pipeline + tests','Structured logging + health','Kubernetes/Helm (designed-for)'])]
    for i, (h, items) in enumerate(cols):
        x = 0.9 + i*3.05
        rect(s, x, 2.0, 2.85, 4.3, fill=BG, line=LINE, radius=0.06)
        rect(s, x, 2.0, 2.85, 0.7, fill=_mix([BLUE,INDIGO,GREEN,AMBER][i], WHITE, 0.12), radius=0.06)
        text(s, x+0.28, 2.16, 2.4, 0.4, [[(h, 15, INK, True)]])
        bullets(s, x+0.28, 2.98, 2.45, 3.2, items, size=11, marker_color=[BLUE,INDIGO,GREEN,AMBER][i], gap=9)
    text(s, 0.9, 6.55, 11.5, 0.4, [[('The reinsurance mathematics lives in a pure, unit-tested domain core — the server orchestrates and persists; it never re-implements the formulas.', 11, SLATE, False)]])
    page_no(s, pn())

# ============================================================== BUSINESS VALUE
def value():
    s = slide(BG); accent_bar(s, color=GREEN); kicker(s, 'Business Value', color=GREEN)
    title(s, 'Where RIOS creates value')
    cards = [('One source of truth','Place, bind, account, reconcile and claims on a single data model — no swivel-chair rekeying between systems.',BLUE),
             ('Correct & reconcilable','Integer-accurate money and a technical → financial chain that reconciles to zero, by construction.',INDIGO),
             ('Audited & inspection-ready','Hash-chained, append-only audit and a compliance surface built in — not bolted on.',GREEN),
             ('Real-time exposure','Capacity utilisation and accumulation monitored live, with breach detection before you bind.',AMBER),
             ('Configurable without code','Statuses, lines of business and rules are metadata — change the business vocabulary with no deployment.',BLUE),
             ('Less system sprawl','Reinsurance, finance, HR, documents and analytics in one platform instead of a dozen.',INDIGO)]
    for i, (h, b, c) in enumerate(cards):
        col = i % 3; row = i // 3
        x = 0.9 + col*4.05; y = 2.0 + row*2.4
        rect(s, x, y, 3.8, 2.2, fill=WHITE, line=LINE, radius=0.07)
        icon_tile(s, x+0.3, y+0.3, color=c, size=0.5)
        text(s, x+0.98, y+0.36, 2.7, 0.5, [[(h, 13.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, x+0.3, y+1.05, 3.25, 1.0, [[(b, 10.5, SLATE, False)]], line_spacing=1.15)
    page_no(s, pn())

# ============================================================== DELIVERED vs DESIGNED-FOR
def maturity():
    s = slide(WHITE); accent_bar(s, color=BLUE); kicker(s, 'Honest Status')
    title(s, 'Delivered today vs designed-for')
    text(s, 0.9, 1.68, 11.5, 0.5, [[('RIOS is a foundation / vertical slice: the reinsurance core is proven end to end; enterprise breadth is architected and named, not pretended. ', 11.5, SLATE, False),
                                     ('Source: the project’s own registers (docs/phases.md, docs/open-questions.md).', 11.5, MUTE, True)]], line_spacing=1.2)
    # delivered column
    rect(s, 0.9, 2.35, 5.75, 4.35, fill=_mix(GREEN, WHITE, 0.07), line=LINE, radius=0.05)
    rect(s, 0.9, 2.35, 5.75, 0.55, fill=GREEN, radius=0.05)
    text(s, 1.15, 2.35, 5.4, 0.55, [[('DELIVERED — built, working, tested', 12.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, 1.15, 3.12, 5.3, 3.4, [
        'Reinsurance core, end to end — integration-tested',
        'Metadata-driven config: 8 structures × 13 lines of business',
        'Underwriting: stage machine, risk score, pricing, cat metrics',
        'Multi-tenant RLS, RBAC, hash-chained audit, integer money',
        'Deterministic grounded assistant — works with AI off',
        'Broker / cedent, capacity, exposure, territory, tasks & SLA',
        'Unit + integration tests; Docker-compose deploy',
    ], size=10.5, color=INK, marker_color=GREEN, gap=10)
    # designed-for column
    rect(s, 6.85, 2.35, 5.75, 4.35, fill=BG, line=LINE, radius=0.05)
    rect(s, 6.85, 2.35, 5.75, 0.55, fill=AMBER, radius=0.05)
    text(s, 7.1, 2.35, 5.4, 0.55, [[('DESIGNED-FOR — architected, not yet built', 12.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, 7.1, 3.08, 5.3, 3.4, [
        'IFRS 17 & Solvency II measurement engines',
        'Full BI: drag-drop report / dashboard designers, pivots, scheduling',
        'Full SSO / SAML, field-level security, managed KMS keys',
        'Portals (broker / cedent / client), native mobile',
        'Microservices, API gateway, Kafka event bus, observability & SLOs',
        'Live connectors (ACORD / bordereaux), API marketplace',
        'Live cat-model APIs, image OCR, optional LLM narration',
    ], size=10, color=SLATE, marker_color=AMBER, gap=5)
    page_no(s, pn())

# ============================================================== PRINCIPLES
def principles():
    s = slide(BG); accent_bar(s, color=INDIGO); kicker(s, 'Design Principles', color=INDIGO)
    title(s, 'The principles behind the build')
    cards = [('Correctness first','Reinsurance maths in a pure, unit-tested core; money never floats.',BLUE),
             ('Metadata-driven','Business vocabulary is data, not code — change it without a release.',INDIGO),
             ('Secure by default','Tenant isolation via RLS, permission-bound routes, least privilege.',GREEN),
             ('Auditable & reconcilable','Every material change hash-chained; the money chain nets to zero.',AMBER),
             ('Grounded AI','Deterministic, explainable, confirmation-gated — never a black box.',BLUE),
             ('Honest about scope','Delivered vs designed-for stated openly — no silent gaps.',INDIGO)]
    for i, (h, b, c) in enumerate(cards):
        col = i % 3; row = i // 3
        x = 0.9 + col*4.05; y = 2.0 + row*2.4
        rect(s, x, y, 3.8, 2.2, fill=WHITE, line=LINE, radius=0.07)
        icon_tile(s, x+0.3, y+0.3, color=c, size=0.5)
        text(s, x+0.98, y+0.36, 2.7, 0.5, [[(h, 13.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, x+0.3, y+1.05, 3.25, 1.0, [[(b, 10.5, SLATE, False)]], line_spacing=1.15)
    page_no(s, pn())

# ============================================================== JOURNEY
def journey():
    s = slide(WHITE); accent_bar(s, color=AMBER); kicker(s, 'Adoption Journey', color=AMBER)
    title(s, 'A guided path from evaluation to go-live')
    steps = ['Discovery','Solution design','Configuration','Data migration','Training','Go-live','Support & roadmap']
    y = 3.4
    for i, t in enumerate(steps):
        x = 0.95 + i*1.72
        c = [BLUE,INDIGO,GREEN,AMBER][i%4]
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
        _noshadow(d); d.fill.solid(); d.fill.fore_color.rgb = c; d.line.fill.background()
        text(s, x-0.35, y-0.02, 1.2, 0.5, [[(str(i+1), 15, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x-0.55, y+0.62, 1.6, 0.7, [[(t, 11, INK, True)]], align=PP_ALIGN.CENTER, line_spacing=1.0)
        if i < len(steps)-1:
            rect(s, x+0.5, y+0.22, 1.22, 0.06, fill=LINE, rounded=False)
    text(s, 0.9, 5.4, 11.5, 0.6, [[('Phased configuration, data migration and role-based training — de-risking the move to a unified platform, with a continuous-improvement roadmap after go-live.', 12.5, SLATE, False)]], line_spacing=1.2)
    page_no(s, pn())

# ============================================================== ROADMAP
def roadmap():
    s = slide(INK); accent_bar(s, color=BLUE)
    text(s, 0.9, 0.7, 8, 0.4, [[('ROADMAP', 11, RGBColor(0x93,0xC5,0xFD), True)]])
    text(s, 0.9, 1.05, 11, 0.8, [[('From foundation to intelligent reinsurance', 28, WHITE, True)]])
    ph = [('Now','Foundation','Unified platform, reinsurance core, audit & security — the delivered slice.',BLUE),
          ('Next','Breadth','Statement lifecycle, bordereaux, richer claims & finance sub-ledgers.',INDIGO),
          ('Later','Regulatory','IFRS 17 & Solvency II engines; governed regulatory report packs.',GREEN),
          ('Future','Connected','Portals, live connectors, BI designers and marketplace.',AMBER),
          ('Vision','Intelligent','Predictive portfolio & capital insight; largely straight-through ops.',BLUE)]
    rect(s, 1.1, 3.35, 11.1, 0.05, fill=RGBColor(0x33,0x41,0x55), rounded=False)
    for i, (yr, h, b, c) in enumerate(ph):
        x = 1.1 + i*2.28
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(3.2), Inches(0.34), Inches(0.34))
        _noshadow(d); d.fill.solid(); d.fill.fore_color.rgb = c; d.line.fill.background()
        text(s, x-0.4, 2.55, 1.2, 0.4, [[(yr, 15, WHITE, True)]], align=PP_ALIGN.CENTER)
        text(s, x-0.55, 3.75, 2.15, 0.5, [[(h, 13, c, True)]], line_spacing=1.0)
        text(s, x-0.55, 4.25, 2.15, 1.6, [[(b, 10.5, RGBColor(0xC7,0xD2,0xFE), False)]], line_spacing=1.15)
    text(s, 0.9, 6.5, 11.5, 0.4, [[('Roadmap shown by horizon, not committed dates — capability direction, honestly staged.', 10.5, MUTE, False)]])
    page_no(s, pn())

# ============================================================== THANK YOU
def thankyou():
    s = slide(INK); accent_bar(s, color=BLUE)
    m = s.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(0.9), Inches(0.85), Inches(0.58), Inches(0.58))
    _noshadow(m); m.fill.solid(); m.fill.fore_color.rgb = BLUE; m.line.fill.background()
    text(s, 1.6, 0.9, 4, 0.5, [[('RIOS', 22, WHITE, True)]])
    text(s, 0.9, 2.7, 8, 1.4, [[('Let’s modernise', 40, WHITE, True)], [('your reinsurance operations.', 40, RGBColor(0x93,0xC5,0xFD), True)]], line_spacing=1.02)
    text(s, 0.95, 4.7, 7, 0.6, [[('Request a live walkthrough and see the unified platform on your own book.', 14, RGBColor(0xC7,0xD2,0xFE), False)]])
    for i, (h, v) in enumerate([('Website','rios.cloud'),('Email','hello@rios.cloud'),('Demo','rios.cloud/demo')]):
        x = 0.95 + i*2.9
        rect(s, x, 5.7, 2.65, 0.95, fill=RGBColor(0x1E,0x29,0x3B), radius=0.1)
        text(s, x+0.25, 5.82, 2.3, 0.3, [[(h.upper(), 9, MUTE, True)]])
        text(s, x+0.25, 6.12, 2.3, 0.4, [[(v, 13, WHITE, True)]])
    browser_frame(s, 'dashboard.png', 9.4, 2.4, 3.4, addr='app.rios.cloud')

# ============================================================== MODULE CONTENT
# (section_key, kicker, title, screenshot, purpose, [capabilities], [inside], [users], ai_line, addr, color)
MODULES = [
  ('OVERVIEW','Overview','Dashboards, Executive & Intelligence','executive.png',
   'The command center. Operational dashboards and executive-intelligence views aggregate live KPIs across every module — written premium, technical result, portfolio mix and pipeline — with drill-down to source, plus enterprise search and a mobile-ready experience.',
   ['Operational dashboard with live tiles and alerts','Executive-intelligence views by leadership persona',
    'Premium, technical-result and portfolio-mix trends','Enterprise search across the whole suite'],
   ['Dashboard','Executive','Intelligence','AI Insights','Search','Mobile'],
   ['Executives','Portfolio managers','Analysts'],
   'AI Insights surfaces grounded, explainable observations per domain — transparent heuristics, not a black box, and fully functional with AI switched off.',
   'app.rios.cloud/executive', BLUE),

  ('UNDERWRITING','Underwriting','Underwriting Workspace','underwriting.png',
   'The underwriting desk. A submission workbench triages risks, applies a transparent risk score and drives referrals and approvals through a guardrailed stage machine — with tasks, collaboration and audit on every action.',
   ['Submission triage, risk scoring and pipeline','Referral routing and a maker-checker approval matrix',
    'Stage machine — illegal transitions rejected','Tasks, collaboration and a full audit trail'],
   ['Submission Workbench','Referrals & Approvals','Analytics','Placement'],
   ['Underwriters','Chief Underwriter','Actuaries'],
   'A deterministic advisor flags missing information, consistency issues and clause suggestions — grounded in the submission, never auto-deciding.',
   'app.rios.cloud/underwriting', INDIGO),

  ('UNDERWRITING','Treaty','Treaty Workspace','treaty.png',
   'The treaty lifecycle in one cockpit: register the treaty, build a priced layer tower, manage versions, clauses and wording, the tax schedule and endorsements, and view the technical account over an audit-stitched timeline.',
   ['Priced layer tower — rate-on-line, reinstatements','Immutable versions and endorsements',
    'Special clauses, wording and tax schedule','Technical account that reconciles to zero'],
   ['Treaty Register','Layer Tower','Clauses & Wording','Endorsements','Technical Account'],
   ['Treaty underwriters','Contract admins','Actuaries'],
   'Grounded clause and data-quality checks assist the drafter; every posting is confirmed — nothing is auto-committed.',
   'app.rios.cloud/w/treaty', INDIGO),

  ('UNDERWRITING','Facultative','Facultative Workspace','facultative.png',
   'A facultative desk. Capture and compare market quotes, build a signed-down placement tower with lead, follow, coinsurance and retro lines, and attach engineering and inspection reports.',
   ['Market quotes with best-quote comparison','Signed-down placement and coinsurance lines',
    'Engineering and inspection report capture','Placement timeline and completeness tracking'],
   ['Quotes','Placement Tower','Coinsurance','Engineering Reports'],
   ['Fac underwriters','Placement brokers'],
   'Deterministic best-quote and placement-gap checks — transparent and explainable, not an opaque model.',
   'app.rios.cloud/w/facultative', INDIGO),

  ('UNDERWRITING','Pricing · Placement · Retro','Pricing, Placement, Retrocession & Adjustments','pricing.png',
   'Actuarial pricing scenarios, placement management, outwards retrocession and treaty adjustments — with the gross / ceded / net position reconciled and linked to recoveries.',
   ['Technical / burning-cost pricing & sensitivity','Placement lines and signed-down orders',
    'Outwards retrocession and recovery linkage','Adjustments: profit commission, sliding scale'],
   ['Pricing','Placement','Retrocession','Adjustments'],
   ['Actuaries','Pricing analysts','Retro managers'],
   'Scenario and sensitivity are computed in the pure domain core — auditable maths, not a black box.',
   'app.rios.cloud/pricing', INDIGO),

  ('UNDERWRITING','Capacity & Exposure','Capacity & Exposure','capacity-exposure.png',
   'Real-time capacity utilisation and exposure accumulation. Monitor limits with red-amber-green status, aggregate exposure by peak zone, and see breach forecasts before you bind.',
   ['Capacity utilisation with RAG status and alerts','Exposure aggregation and peak-zone concentration',
    'Peril × geography accumulation heatmap','Straight-line breach forecast prior to binding'],
   ['Capacity','Exposure Management','Accumulation'],
   ['Exposure managers','Cat analysts','Underwriters'],
   'Rule-based breach and concentration alerts — grounded in live limits, every factor explainable.',
   'app.rios.cloud/w/capacity-exposure', INDIGO),

  ('UNDERWRITING','Territory','Territory Workspace','territory.png',
   'The geographic master. A country → zone hierarchy joined to CRESTA / peril / risk zones, each linked to live exposure with total insured value, modelled PML and a transparent risk score.',
   ['Country → zone hierarchy','CRESTA / peril / risk zones',
    'Insured value and modelled PML per territory','Transparent, factor-based zone risk score'],
   ['Territory Hierarchy','Risk Zones','Exposure Link'],
   ['Exposure managers','Cat modellers'],
   'Transparent zone risk scoring from your own exposure — every contributing factor is visible.',
   'app.rios.cloud/w/territory', INDIGO),

  ('DISTRIBUTION','Distribution','Parties, Brokers, Cedents & CRM','parties.png',
   'A single counterparty backbone. Parties, clients, brokers and cedents each carry a 360° profile — contracts, statements and claims — and a CRM manages pipeline, opportunities and communications.',
   ['Party 360: contracts, statements, claims','Broker and cedent profiles, tiers and performance',
    'CRM pipeline, opportunities and activities','Contact directory and communication log'],
   ['Parties','Clients','Brokers','Cedents','CRM'],
   ['Relationship managers','Broker managers','Business development'],
   'A transparent, unit-tested relationship & profitability score highlights where to focus — no black box.',
   'app.rios.cloud/parties', GREEN),

  ('OPERATIONS','Claims','Claims, Bordereaux & Recoveries','claims.png',
   'The claims flow — notify → reserve movement → paid loss — with bordereaux views, reserves and recoveries, connected to treaties, finance and the audit log.',
   ['Claim flow with reserve movements and payments','Bordereaux views and reconciliation',
    'Recovery and cash-call tracking','Every movement audited and reconcilable'],
   ['Claims','Bordereaux','Recoveries'],
   ['Claims handlers','Adjusters','Recoveries team'],
   'Deterministic reserve-consistency and triage checks assist adjusters; controls are never bypassed.',
   'app.rios.cloud/claims', AMBER),

  ('OPERATIONS','Operations & Workflow','Operations Center, Workflow & Audit','workflow-engine.png',
   'The operational control tower. Live workflow instances, SLA-scored tasks, an escalation queue and an approval matrix — over a tamper-evident, hash-chained audit trail.',
   ['Operations center with live work queues','SLA scoring with tiered escalation',
    'Approval matrix and delegation','Hash-chained, append-only audit log'],
   ['Operations Center','Workflow Center','Audit Log'],
   ['Operations','Team leads','Auditors'],
   'Rule-based SLA and bottleneck signals over the live work queues — explainable, not predictive magic.',
   'app.rios.cloud/workflow-engine', AMBER),

  ('FINANCE','Technical & General Accounting','Technical & General Accounting','accounting.png',
   'The accounting core. Reconcilable financial events post balanced general-ledger entries; technical accounts and statements tie premium, commission and claims to the ledger, with a period-close checklist and procurement.',
   ['Balanced GL postings from financial events','Technical accounts and statements of account',
    'Period-close checklist','Procurement and payables'],
   ['Accounting','Statements','Period Close','Procurement'],
   ['Accountants','Technical accountants','Controllers'],
   'Grounded posting and reconciliation checks; the assistant proposes, you confirm — money is integer-exact.',
   'app.rios.cloud/accounting', BLUE),

  ('FINANCE','Treasury & Investment','Treasury & Investment','finance.png',
   'Treasury and investment views over the GL core. Cash and bank positions, settlements and an investment register — with a consolidated finance workspace. (Full sub-ledgers are designed-for.)',
   ['GL-linked cash and bank views','Investment register and positions',
    'Consolidated finance workspace','Integer-accurate money — no floating point'],
   ['Finance','Treasury','Investments'],
   ['Treasury','Finance ops','Investment ops'],
   'Money is stored in integer minor units for exactness; liquidity signals are transparent and auditable.',
   'app.rios.cloud/finance', BLUE),

  ('ANALYTICS','Reporting & Analytics','Reporting & Analytics','reports.png',
   'Reporting over the suite. A report library with CSV and print-to-PDF export, scheduled report definitions, and portfolio analytics with drill-down. (Drag-drop designers and pivots are designed-for.)',
   ['Report library with CSV / print-to-PDF export','Scheduled report definitions',
    'Portfolio analytics and KPI dashboards','Drill-down from summary to source'],
   ['Reports','Scheduled Reports','Analytics'],
   ['Analysts','Executives','Actuaries'],
   'Executive summaries are template-merged over live KPIs; an optional LLM may narrate — it never invents the numbers.',
   'app.rios.cloud/reports', INDIGO),

  ('ANALYTICS','Risk & Capital','Risk & Capital','risk-capital.png',
   'Risk and capital views. Catastrophe metrics (AAL / PML / EP / TVaR) via a cat-model adapter, risk-appetite monitoring, reserving and loss-development. (Full IFRS 17 / Solvency II engines are designed-for.)',
   ['Catastrophe metrics — AAL / PML / EP / TVaR','Risk-appetite and capital views',
    'Reserving & loss development (chain-ladder IBNR)','Scenario analysis'],
   ['Risk & Capital','Reserving','Scenarios'],
   ['Risk officers','Actuaries','Capital team'],
   'Scenario maths runs in the pure domain core — transparent, unit-tested and auditable.',
   'app.rios.cloud/risk-capital', INDIGO),

  ('ANALYTICS','Regulatory & Compliance','Regulatory, Compliance & Returns','compliance.png',
   'A single assurance surface. An audit dashboard verifies chain integrity; approvals, activity and data-access logs, a compliance calendar and regulatory return workspaces keep the reinsurer inspection-ready.',
   ['Audit dashboard and chain-integrity check','Approvals, activity and data-access logs',
    'Compliance calendar with due dates','Regulatory return workspaces'],
   ['Regulatory','Compliance','Returns','Audit'],
   ['Compliance','Regulatory reporting','Auditors'],
   'The audit chain is cryptographically verified; compliance signals are rule-based and explainable.',
   'app.rios.cloud/compliance', INDIGO),

  ('HRMS','Human Resources','Human Resources — People, Payroll & Org','attendance.png',
   'An integrated HR back office. People records with an audited status lifecycle, an attendance command center with manager approvals routed by the real org hierarchy, payroll, performance, assets and org structure.',
   ['People records with audited status lifecycle','Attendance command center; OD / WFH / regularization',
    'Manager approvals via the org hierarchy','Payroll, performance, assets and org structure'],
   ['Attendance','People','Payroll','Performance','Assets','Org Structure'],
   ['HR','People managers','Payroll'],
   'Rule-based attendance-anomaly checks; approvals route to the manager resolved from the real org hierarchy.',
   'app.rios.cloud/attendance', GREEN),

  ('PLATFORM','Master Data','Master Data & Products','products.png',
   'The governed reference layer. Products and coverage definitions, lines of business, currencies, countries, clauses and code-lists are configured as metadata — new values are added without a deployment.',
   ['Products and coverage definitions','Lines of business, currencies, countries',
    'Code-lists as metadata — no hard-coded enums','New values added without a release'],
   ['Products','Reference Data','Code Lists','Business Rules'],
   ['Data stewards','Config admins'],
   'Deterministic validation keeps reference data clean — this is configuration, not code.',
   'app.rios.cloud/products', AMBER),

  ('PLATFORM','Documents','Documents & Knowledge','documents.png',
   'A document hub. A repository with versioning and approval routing, template and clause libraries, retention policies and a searchable knowledge base. (Blob storage and image OCR are wired at deployment.)',
   ['Repository with versioning and supersede chains','Templates and clause / wording libraries',
    'Approval routing and retention','Knowledge base and standard procedures'],
   ['Documents','Templates','Knowledge Base'],
   ['All users','Knowledge managers'],
   'Deterministic field extraction from text is built in; the image → text step uses an external OCR engine.',
   'app.rios.cloud/documents', AMBER),

  ('PLATFORM','Integration','Integration & Automation','integration-hub.png',
   'Open by design. REST / GraphQL APIs and webhooks, a transactional outbox + relay, a connector registry with one-time API keys, messaging, and an automation studio composing rules and events. (Live sinks wired per deployment.)',
   ['REST / GraphQL APIs and webhooks','Transactional outbox + relay; connector registry',
    'Messaging queue (email / SMS) — sink at deployment','Automation studio, scheduler and portal'],
   ['Integration Hub','Messaging','Automation Studio','Portal'],
   ['IT','Integration engineers','Ops'],
   'Config validation and mapping assistance; the outbox + relay move events reliably and transactionally.',
   'app.rios.cloud/integration-hub', AMBER),

  ('ADMIN','Administration','Administration & Security','admin.png',
   'Everything administrators need. Users, roles and permissions; legal entities and delegation of authority; security operations, retention, cost management and feature flags. (Field-level security is designed-for.)',
   ['Roles and a permission vocabulary','OIDC SSO, MFA and security operations',
    'Legal entities and delegation of authority','Retention, cost management and feature flags'],
   ['Admin','Legal Entities','Ops Console','Delegation','Security','Security Ops','Retention','Cost Mgmt','Features'],
   ['Administrators','Security','IT'],
   'Rule-based access and policy checks; the guardrailed assistant re-checks permissions before any action.',
   'app.rios.cloud/admin', BLUE),
]

sec = {
  'OVERVIEW':      ('01','Overview & Intelligence','Command center — dashboards, executive KPIs, AI insights, search and mobile.', BLUE),
  'UNDERWRITING':  ('02','Underwriting — the Reinsurance Core','Treaty & facultative, placement, pricing, capacity, exposure, territory, retrocession and adjustments.', INDIGO),
  'DISTRIBUTION':  ('03','Distribution','Parties, clients, brokers, cedents and the CRM.', GREEN),
  'OPERATIONS':    ('04','Operations','Claims, bordereaux, recoveries, the operations & workflow control tower and audit.', AMBER),
  'FINANCE':       ('05','Finance & Accounting','Technical and general accounting, statements, treasury and investment, period close and procurement.', BLUE),
  'ANALYTICS':     ('06','Analytics & Compliance','Reporting, analytics, risk & capital, regulatory returns and compliance.', INDIGO),
  'HRMS':          ('07','Human Resources','People, attendance, payroll, performance, assets and organisation structure.', GREEN),
  'PLATFORM':      ('08','Master Data, Documents & Integration','The governed reference layer, the document hub and open integration & automation.', AMBER),
  'ADMIN':         ('09','Administration & Security','Users and roles, legal entities, delegation, security operations, retention, cost and features.', BLUE),
}

# ============================================================== ASSEMBLE
cover()
about()
problems()
product_map()
architecture()
security()

seen = set()
for m in MODULES:
    key = m[0]
    if key not in seen:
        seen.add(key)
        num, ttl, sub, col = sec[key]
        section(num, ttl, sub, col)
    module_slide(*m[1:])

assistant()
lifecycle()
technology()
value()
maturity()
principles()
journey()
roadmap()
thankyou()

out = os.path.join(BASE, 'RIOS_Enterprise_Presentation.pptx')
prs.save(out)
print('saved', out, '·', len(prs.slides._sldIdLst), 'slides')
