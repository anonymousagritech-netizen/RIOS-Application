#!/usr/bin/env python3
"""Approximate pptx -> PNG previewer (PIL). Not pixel-perfect, but faithful enough
to judge layout, spacing, overlap and colour. Renders each slide + a contact sheet."""
import os, io, sys
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw, ImageFont

BASE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(BASE, 'RIOS_Enterprise_Presentation.pptx')
OUT = os.path.join(BASE, 'preview'); os.makedirs(OUT, exist_ok=True)
REG = '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf'
BLD = '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf'
SCALE = 150  # px per inch
EMU_IN = 914400
_fc = {}
def font(size_pt, bold):
    key = (round(size_pt), bold)
    if key not in _fc:
        px = max(6, int(size_pt * SCALE / 72))
        _fc[key] = ImageFont.truetype(BLD if bold else REG, px)
    return _fc[key]

def emu_px(v):
    return int((v / EMU_IN) * SCALE)

def rgb(c):
    return (c[0], c[1], c[2])

def rounded(dr, box, radius, fill, outline, width=1):
    x0, y0, x1, y1 = box
    if x1 <= x0 or y1 <= y0:
        return
    r = max(0, min(radius, (x1 - x0) // 2, (y1 - y0) // 2))
    if r > 0:
        dr.rounded_rectangle(box, radius=r, fill=fill, outline=outline, width=width)
    else:
        dr.rectangle(box, fill=fill, outline=outline, width=width)

def shape_fill(sh):
    try:
        f = sh.fill
        if f.type is not None and f.type == 1:  # solid
            return rgb(f.fore_color.rgb)
    except Exception:
        pass
    return None

def shape_line(sh):
    try:
        ln = sh.line
        if ln.color and ln.color.type is not None:
            return rgb(ln.color.rgb)
    except Exception:
        pass
    return None

def draw_text(dr, sh):
    tf = sh.text_frame
    L = emu_px(sh.left); T = emu_px(sh.top); W = emu_px(sh.width); H = emu_px(sh.height)
    wrap = getattr(tf, 'word_wrap', True) is not False
    # build wrapped lines: each line is a list of (word, size, color, bold)
    lines = []
    for p in tf.paragraphs:
        toks = []
        for r in p.runs:
            if not r.text:
                continue
            s = r.font.size.pt if r.font.size else 12
            col = rgb(r.font.color.rgb) if (r.font.color and r.font.color.type is not None) else (30,41,59)
            b = bool(r.font.bold)
            parts = r.text.split(' ')
            for k, w in enumerate(parts):
                toks.append((w + (' ' if k < len(parts)-1 else ''), s, col, b))
        align = p.alignment
        if not toks:
            lines.append(([], align)); continue
        cur = []; cw = 0
        for tok in toks:
            tw = dr.textlength(tok[0], font=font(tok[1], tok[3]))
            if wrap and cur and cw + tw > W:
                lines.append((cur, align)); cur = []; cw = 0
            cur.append(tok); cw += tw
        if cur:
            lines.append((cur, align))
    line_h = []
    for toks, _ in lines:
        mh = max([font(s, b).size for (_, s, _, b) in toks], default=int(12*SCALE/72))
        line_h.append(int(mh * 1.32) if toks else int(9*SCALE/72))
    total = sum(line_h)
    anchor = tf.vertical_anchor
    y = T
    if anchor == MSO_ANCHOR.MIDDLE:
        y = T + max(0, (H - total)//2)
    elif anchor == MSO_ANCHOR.BOTTOM:
        y = T + max(0, H - total)
    for (toks, align), lh in zip(lines, line_h):
        if not toks:
            y += lh; continue
        widths = [dr.textlength(t, font=font(s, b)) for (t, s, _, b) in toks]
        tw = sum(widths)
        if align == PP_ALIGN.CENTER:
            x = L + max(0, (W - tw)//2)
        elif align == PP_ALIGN.RIGHT:
            x = L + max(0, W - tw)
        else:
            x = L
        for (t, s, col, b), w in zip(toks, widths):
            dr.text((x, y), t, font=font(s, b), fill=col)
            x += w
        y += lh

def render_slide(slide, idx):
    W = emu_px(prs.slide_width); H = emu_px(prs.slide_height)
    img = Image.new('RGB', (W, H), (255, 255, 255))
    dr = ImageDraw.Draw(img)
    for sh in slide.shapes:
        try:
            st = sh.shape_type
            L = emu_px(sh.left); T = emu_px(sh.top); Wd = emu_px(sh.width); Hd = emu_px(sh.height)
            box = (L, T, L+Wd, T+Hd)
            if st == MSO_SHAPE_TYPE.PICTURE:
                im = Image.open(io.BytesIO(sh.image.blob)).convert('RGB')
                im = im.resize((max(1,Wd), max(1,Hd)))
                img.paste(im, (L, T))
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.CHART or sh.has_chart:
                rounded(dr, box, 6, (238,244,255), (200,210,230), 1)
                dr.text((L+8, T+8), '[chart]', font=font(12, True), fill=(80,90,120))
                continue
            # autoshape / textbox
            fill = shape_fill(sh)
            line = shape_line(sh)
            sname = ''
            try:
                sname = str(sh.auto_shape_type)
            except Exception:
                pass
            if fill or line:
                if 'OVAL' in sname:
                    dr.ellipse(box, fill=fill, outline=line, width=2 if line else 1)
                elif 'DOWN_ARROW' in sname:
                    cx = L + Wd/2; sw = Wd*0.42/2; hh = Hd*0.5
                    dr.polygon([(cx-sw,T),(cx+sw,T),(cx+sw,T+Hd-hh),(L+Wd,T+Hd-hh),
                                (cx,T+Hd),(L,T+Hd-hh),(cx-sw,T+Hd-hh)], fill=fill or (150,160,175))
                elif 'RIGHT_ARROW' in sname:
                    cy = T + Hd/2; sh2 = Hd*0.44/2; hw = Wd*0.5
                    dr.polygon([(L,cy-sh2),(L+Wd-hw,cy-sh2),(L+Wd-hw,T),(L+Wd,cy),
                                (L+Wd-hw,T+Hd),(L+Wd-hw,cy+sh2),(L,cy+sh2)], fill=fill or (150,160,175))
                elif 'HEXAGON' in sname:
                    q = Wd*0.25
                    dr.polygon([(L+q,T),(L+Wd-q,T),(L+Wd,T+Hd/2),(L+Wd-q,T+Hd),
                                (L+q,T+Hd),(L,T+Hd/2)], fill=fill, outline=line)
                else:
                    rounded(dr, box, guess_radius(sh, Wd, Hd), fill, line, 2 if line else 1)
            if sh.has_text_frame and sh.text_frame.text.strip():
                draw_text(dr, sh)
        except Exception as e:
            pass
    path = os.path.join(OUT, f'slide_{idx:02d}.png')
    img.save(path)
    return img

def guess_radius(sh, w, h):
    try:
        if 'ROUNDED' in str(sh.auto_shape_type):
            return max(6, int(min(w, h) * 0.12))
    except Exception:
        pass
    return 0

def radius_fill(sh):
    return None

prs = Presentation(SRC)
imgs = []
only = None
if len(sys.argv) > 1:
    only = set(int(x) for x in sys.argv[1:])
for i, s in enumerate(prs.slides, 1):
    if only and i not in only:
        continue
    imgs.append(render_slide(s, i))
print('rendered', len(imgs), 'slides to', OUT)

# contact sheet (6 cols)
if imgs and not only:
    cols = 6
    tw, th = imgs[0].size
    sc = 360 / tw
    cw, ch = int(tw*sc), int(th*sc)
    rows = (len(imgs)+cols-1)//cols
    sheet = Image.new('RGB', (cols*cw + (cols+1)*10, rows*ch + (rows+1)*10), (235,238,242))
    for i, im in enumerate(imgs):
        r, c = divmod(i, cols)
        t = im.resize((cw, ch))
        sheet.paste(t, (10 + c*(cw+10), 10 + r*(ch+10)))
    sheet.save(os.path.join(OUT, 'contact_sheet.png'))
    print('contact sheet saved')
