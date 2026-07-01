#!/usr/bin/env python3
"""Generate RIOS_Enterprise_Presentation.pptx — a premium, light-theme executive
sales deck for RIOS, embedding real product screenshots inside browser mockups."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from PIL import Image

BASE = os.path.dirname(os.path.abspath(__file__))
SCR = os.path.join(BASE, 'assets', 'screens')

# ---- palette ----
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG    = RGBColor(0xF8, 0xFA, 0xFC)
PALE  = RGBColor(0xEE, 0xF4, 0xFF)
BLUE  = RGBColor(0x25, 0x63, 0xEB)
INDIGO= RGBColor(0x4F, 0x46, 0xE5)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
INK   = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x47, 0x55, 0x69)
MUTE  = RGBColor(0x94, 0xA3, 0xB8)
LINE  = RGBColor(0xE2, 0xE8, 0xF0)
FONT  = 'Segoe UI'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s


def _noshadow(sp):
    sp.shadow.inherit = False


def rect(s, x, y, w, h, fill=None, line=None, rounded=True, radius=0.08, line_w=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    _noshadow(shp)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(sp_after); p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=13, color=SLATE, gap=6, marker_color=BLUE):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        rm = p.add_run(); rm.text = '▸  '
        rm.font.size = Pt(size); rm.font.color.rgb = marker_color; rm.font.bold = True; rm.font.name = FONT
        r = p.add_run(); r.text = it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def chip(s, x, y, w, txt, fill=PALE, fg=BLUE, h=0.42, size=11, bold=True):
    c = rect(s, x, y, w, h, fill=fill, radius=0.5)
    text(s, x, y - 0.02, w, h, [[(txt, size, fg, bold)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c


def icon_tile(s, x, y, glyph, color=BLUE, size=0.62, gsize=22):
    t = rect(s, x, y, size, size, fill=None, radius=0.28)
    t.fill.solid(); t.fill.fore_color.rgb = _mix(color, WHITE, 0.14)
    text(s, x, y, size, size, [[(glyph, gsize, color, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return t


def _mix(a, b, t):
    return RGBColor(int(a[0]+(b[0]-a[0])*(1-t)) if False else int(a[0]*(t)+b[0]*(1-t)),
                    int(a[1]*t+b[1]*(1-t)), int(a[2]*t+b[2]*(1-t)))


def accent_bar(s, x=0.0, y=0.0, w=0.18, h=7.5, color=BLUE):
    rect(s, x, y, w, h, fill=color, rounded=False)


def kicker(s, txt, x=0.9, y=0.6, color=BLUE):
    text(s, x, y, 8, 0.3, [[(txt.upper(), 11, color, True)]])


def title(s, txt, x=0.9, y=0.9, size=30, w=11.5, color=INK):
    text(s, x, y, w, 1.0, [[(txt, size, color, True)]])


def page_no(s, n):
    text(s, 12.3, 7.05, 0.9, 0.3, [[(str(n), 9, MUTE, False)]], align=PP_ALIGN.RIGHT)
    text(s, 0.9, 7.05, 4, 0.3, [[('RIOS · Reinsurance Intelligence & Operations Suite', 9, MUTE, False)]])


def browser_frame(s, img, x, y, w, addr='app.rios.cloud'):
    """Place a screenshot inside a light browser mockup; crop to frame aspect."""
    header_h = 0.34
    body_h = w * 0.60  # 5:3 body
    total_h = header_h + body_h
    # shadow card
    rect(s, x-0.06, y-0.04, w+0.12, total_h+0.12, fill=RGBColor(0xED,0xF1,0xF7), radius=0.05)
    # frame
    rect(s, x, y, w, total_h, fill=WHITE, line=LINE, radius=0.05, line_w=1.0)
    # header bar
    hdr = rect(s, x, y, w, header_h, fill=RGBColor(0xF1,0xF5,0xF9), line=None, radius=0.05)
    for i, col in enumerate([RGBColor(0xFB,0x71,0x85), RGBColor(0xFB,0xBF,0x24), RGBColor(0x34,0xD3,0x99)]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.18+i*0.22), Inches(y+0.12), Inches(0.11), Inches(0.11))
        _noshadow(d); d.fill.solid(); d.fill.fore_color.rgb = col; d.line.fill.background()
    rect(s, x+1.0, y+0.07, w-1.6, 0.2, fill=WHITE, line=LINE, radius=0.5, line_w=0.75)
    text(s, x+1.15, y+0.055, w-1.8, 0.2, [[(addr, 9, MUTE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    # image (crop to body aspect)
    path = os.path.join(SCR, img)
    if os.path.exists(path):
        im = Image.open(path); iw, ih = im.size
        target = (w) / body_h
        cw, ch = iw, int(iw/target)
        if ch > ih:
            ch = ih; cw = int(ih*target)
        im2 = im.crop((0, 0, cw, ch))
        tmp = os.path.join(SCR, '_c_'+img)
        im2.save(tmp)
        s.shapes.add_picture(tmp, Inches(x+0.03), Inches(y+header_h), Inches(w-0.06), Inches(body_h-0.03))
    else:
        rect(s, x+0.03, y+header_h, w-0.06, body_h-0.03, fill=PALE, radius=0.02)
        text(s, x, y+header_h, w, body_h, [[('RIOS · '+img.replace('.png',''), 16, BLUE, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return total_h


def arrow(s, x, y, w, h, color=MUTE):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    _noshadow(a); a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()
    return a
